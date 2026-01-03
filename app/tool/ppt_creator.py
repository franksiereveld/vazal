import os
import requests
import io
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from app.tool.base import BaseTool, ToolResult

class PPTCreatorTool(BaseTool):
    name: str = "ppt_creator"
    description: str = """
    Creates a PowerPoint presentation with text, images, quotes, and sources.
    * Supports using a custom template (.pptx or .potx).
    * Can add multiple slides with titles, bullet points, images, quotes, and sources.
    * AUTOMATICALLY downloads images if you provide a URL (http/https).
    * Saves downloaded images to 'input/' and the final PPT to 'output/'.
    """
    parameters: dict = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The name of the output file (e.g., 'presentation.pptx'). Will be saved in 'output/' folder.",
            },
            "template": {
                "type": "string",
                "description": "Optional path to a template file (.pptx or .potx).",
            },
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "List of bullet points",
                        },
                        "image_path": {
                            "type": "string",
                            "description": "Local path OR URL (http/https) to an image.",
                        },
                        "quote": {
                            "type": "string",
                            "description": "A relevant quote to display if no image is available (Fallback).",
                        },
                        "quote_source": {
                            "type": "string",
                            "description": "The source of the quote (Person, Role, Date). Required if quote is provided.",
                        },
                        "sources": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "List of sources/citations for this slide.",
                        },
                    },
                    "required": ["title", "content"],
                },
            },
        },
        "required": ["filename", "slides"],
    }

    def _download_image(self, url: str) -> str:
        """
        Downloads an image from a URL to the 'input/' folder.
        Returns the local file path.
        """
        try:
            # Create input directory if it doesn't exist
            input_dir = "input"
            if not os.path.exists(input_dir):
                os.makedirs(input_dir)

            # Extract filename from URL or generate a random one
            filename = url.split("/")[-1].split("?")[0]
            if not filename or "." not in filename:
                import uuid
                filename = f"image_{uuid.uuid4().hex[:8]}.jpg"
            
            local_path = os.path.join(input_dir, filename)

            # Download the file with User-Agent to avoid 403 Forbidden
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            
            # Check Content-Type
            content_type = response.headers.get('Content-Type', '').lower()
            if 'image' not in content_type:
                print(f"⚠️ Skipping non-image URL {url} (Content-Type: {content_type})")
                return None

            # Verify image integrity
            try:
                image_data = response.content
                img = Image.open(io.BytesIO(image_data))
                img.verify()
            except Exception as e:
                print(f"⚠️ Invalid image data from {url}: {e}")
                return None
            
            with open(local_path, 'wb') as f:
                f.write(image_data)
            
            print(f"⬇️ Downloaded image: {url} -> {local_path}")
            return local_path
        except Exception as e:
            print(f"⚠️ Failed to download image {url}: {e}")
            return None

    def _find_best_layout(self, prs):
        """
        Finds the best layout that has both a body text placeholder and a picture placeholder.
        Prioritizes Layout 0 for the specific user template.
        """
        # 1. Try to find a layout with specific structure (Body + Picture)
        for i, layout in enumerate(prs.slide_layouts):
            has_body = False
            has_picture = False
            for shape in layout.placeholders:
                # Check for Body (2) or Object (7)
                if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    has_body = True
                # Check for Picture (18)
                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    has_picture = True
            
            if has_body and has_picture:
                print(f"✅ Found optimal layout: Index {i} ({layout.name})")
                return layout

        # 2. Fallback: If using the specific user template, we know Layout 0 is the one.
        # We can check if Layout 0 has at least 2 placeholders (Title + Content)
        if len(prs.slide_layouts) > 0:
             # Heuristic: Layout 0 is often Title Slide, but in this custom template it's the main one.
             # Let's check if Layout 0 has a body placeholder.
             layout_0 = prs.slide_layouts[0]
             for shape in layout_0.placeholders:
                 if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                     print(f"ℹ️ Using Layout 0 ({layout_0.name}) as fallback.")
                     return layout_0

        # 3. Last Resort: Standard "Title and Content" is usually index 1
        if len(prs.slide_layouts) > 1:
            print("⚠️ Using standard Layout 1 fallback.")
            return prs.slide_layouts[1]
        
        return prs.slide_layouts[0]

    async def execute(self, filename: str, slides: list, template: str = None, **kwargs) -> ToolResult:
        # --- STRICT VALIDATION ---
        errors = []
        
        # 1. Check for Minimum Slide Count
        if len(slides) < 5:
            errors.append(f"Presentation is too short ({len(slides)} slides). It MUST have at least 5 slides (Title + Summary + 2 Content + Conclusion).")

        # 2. Check for Executive Summary (Slide 2)
        if len(slides) > 1:
            slide2_title = slides[1].get("title", "").lower()
            valid_summary_terms = ["summary", "agenda", "table of contents", "overview", "roadmap", "table des matières", "sommaire", "inhaltsverzeichnis", "zusammenfassung", "überblick", "agenda"]
            if not any(term in slide2_title for term in valid_summary_terms):
                errors.append("Slide 2 MUST be an 'Executive Summary', 'Agenda', or 'Table of Contents'.")

        # 3. Check for Conclusion (Last Slide)
        if len(slides) > 0:
            last_slide_title = slides[-1].get("title", "").lower()
            valid_conclusion_terms = ["conclusion", "next steps", "future outlook", "summary", "closing", "fazit", "ausblick", "zusammenfassung", "prochaines étapes", "conclusion"]
            if not any(term in last_slide_title for term in valid_conclusion_terms):
                errors.append("The LAST slide MUST be a 'Conclusion', 'Next Steps', or 'Future Outlook'.")

        # 4. Check for Unique Images (if provided)
        image_urls = [s.get("image_path") for s in slides if s.get("image_path")]
        if len(image_urls) != len(set(image_urls)):
            errors.append("You reused the same image on multiple slides. Each slide MUST have a UNIQUE image.")

        # 5. Check for Titles and Content Depth
        for i, slide in enumerate(slides):
            if not slide.get("title"):
                errors.append(f"Slide {i+1} is missing a title.")
            
            # Skip title slide (index 0) and summary (index 1) and conclusion (last index) for strict bullet count check
            # BUT ensure they are not empty.
            content = slide.get("content", [])
            
            if i == 0: # Title Slide
                pass 
            elif i == 1: # Summary Slide
                if len(content) < 3:
                     errors.append(f"Slide 2 (Summary) needs at least 3 items.")
            elif i == len(slides) - 1: # Conclusion Slide
                if len(content) < 2:
                     errors.append(f"Last Slide (Conclusion) needs at least 2 items.")
            else: # Content Slides
                if len(content) < 4:
                    errors.append(f"Slide {i+1} ('{slide.get('title')}') has only {len(content)} bullet points. It needs at least 4 detailed points.")

        if errors:
            error_msg = "❌ PRESENTATION REJECTED due to quality rules:\n" + "\n".join(errors)
            return ToolResult(error=error_msg)
        # -------------------------

        try:
            # Ensure output directory exists
            output_dir = "output"
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # Prepend output directory to filename
            output_path = os.path.join(output_dir, filename)

            # Load template or create new
            if template:
                if not os.path.exists(template):
                    return ToolResult(error=f"Template file not found: {template}")
                prs = Presentation(template)
                print(f"✅ Using template: {template}")
            else:
                # Check for default template 'PPT.pptx'
                default_template = "/home/ubuntu/upload/PPT.pptx"
                if os.path.exists(default_template):
                    prs = Presentation(default_template)
                    print(f"✅ Using default template: {default_template}")
                elif os.path.exists("PPT.pptx"):
                    prs = Presentation("PPT.pptx")
                    print(f"✅ Using default template: PPT.pptx")
                else:
                    prs = Presentation()
                    print("ℹ️ No template provided and default 'PPT.pptx' not found. Using blank theme.")

            # NOTE: Removed slide clearing logic to prevent file corruption.
            # We will append new slides to the existing ones.
            # If the user wants a clean start, they should provide a clean template.
            
            # Find the best layout for content slides
            content_layout = self._find_best_layout(prs)
            
            all_sources = []

            for i, slide_data in enumerate(slides):
                # Use Title Slide layout (usually 0) for the first slide if it's a standard template
                # BUT for this specific user template, Layout 0 is the content layout.
                
                slide = prs.slides.add_slide(content_layout)

                # Set Title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "Untitled")

                # --- Set Content (Bullet Points) ---
                body_shape = None
                
                # Priority 1: Look for specific index 10 (User Template)
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 10:
                        body_shape = shape
                        break
                
                # Priority 2: Look for standard Body (2) or Object (7)
                if not body_shape:
                    for shape in slide.placeholders:
                        if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                            body_shape = shape
                            break
                
                # Priority 3: Any text frame that isn't title
                if not body_shape:
                    for shape in slide.placeholders:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            body_shape = shape
                            break

                if body_shape:
                    tf = body_shape.text_frame
                    tf.text = ""  # Clear default text
                    for point in slide_data.get("content", []):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                else:
                    # Fallback: Create a text box manually
                    print("ℹ️ No body placeholder found. Creating manual text box.")
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(5.0)
                    height = Inches(5.0)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for point in slide_data.get("content", []):
                        p = tf.add_paragraph()
                        p.text = f"• {point}"
                        p.level = 0

                # --- Add Image OR Quote ---
                image_path = slide_data.get("image_path")
                quote = slide_data.get("quote")
                quote_source = slide_data.get("quote_source")
                image_added = False

                if image_path:
                    # Check if it's a URL
                    if image_path.startswith("http://") or image_path.startswith("https://"):
                        image_path = self._download_image(image_path)

                    # Check if local file exists
                    if image_path and os.path.exists(image_path):
                        try:
                            pic_placeholder = None
                            
                            # Priority 1: Look for specific index 11 (User Template)
                            for shape in slide.placeholders:
                                if shape.placeholder_format.idx == 11:
                                    pic_placeholder = shape
                                    break
                            
                            # Priority 2: Look for Picture (18)
                            if not pic_placeholder:
                                for shape in slide.placeholders:
                                    if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                                        pic_placeholder = shape
                                        break
                            
                            if pic_placeholder:
                                # Robust insertion: Get coordinates, then add new picture
                                try:
                                    # REMOVE PLACEHOLDER before inserting picture to avoid "Click to add text" ghosting
                                    # Actually, insert_picture replaces the placeholder content, but sometimes the placeholder remains.
                                    # The best way is to use the placeholder's insert_picture method if available.
                                    pic_placeholder.insert_picture(image_path)
                                    image_added = True
                                except AttributeError:
                                    # Fallback for generic placeholders
                                    left = pic_placeholder.left
                                    top = pic_placeholder.top
                                    width = pic_placeholder.width
                                    height = pic_placeholder.height
                                    
                                    # Remove the placeholder shape itself to avoid "Click to add text"
                                    sp = pic_placeholder.element
                                    sp.getparent().remove(sp)
                                    
                                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                                    image_added = True
                            else:
                                # Fallback: Add image to the right side
                                print("ℹ️ No picture placeholder found. Adding manual image.")
                                left = Inches(5.5)
                                top = Inches(2)
                                height = Inches(3.5)
                                slide.shapes.add_picture(image_path, left, top, height=height)
                                image_added = True
                        except Exception as e:
                            print(f"⚠️ Error adding image {image_path}: {e}")
                    else:
                        print(f"⚠️ WARNING: Skipping invalid or missing image path: {image_path}")

                # --- Fallback: Add Quote if no image ---
                if not image_added and quote:
                    print(f"ℹ️ Adding quote fallback for slide '{slide_data.get('title')}'")
                    # Try to find the picture placeholder to put the quote in
                    quote_placeholder = None
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 11: # Right side
                            quote_placeholder = shape
                            break
                    
                    if quote_placeholder:
                        left = quote_placeholder.left
                        top = quote_placeholder.top
                        width = quote_placeholder.width
                        height = quote_placeholder.height
                        
                        # CRITICAL FIX: Remove the placeholder to get rid of "Click to add text"
                        sp = quote_placeholder.element
                        sp.getparent().remove(sp)
                        
                    else:
                        left = Inches(5.5)
                        top = Inches(2)
                        width = Inches(4.0)
                        height = Inches(3.5)
                    
                    # Create text box for quote
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    # Add Quote Text
                    p = tf.add_paragraph()
                    p.text = f'"{quote}"'
                    p.font.size = Pt(24)
                    p.font.italic = True
                    p.font.color.rgb = RGBColor(80, 80, 80) # Dark gray
                    
                    # Add Quote Source (if provided)
                    if quote_source:
                        p_source = tf.add_paragraph()
                        p_source.text = f"- {quote_source}"
                        p_source.font.size = Pt(14) # Smaller font
                        p_source.font.italic = False # Normal style
                        p_source.font.color.rgb = RGBColor(100, 100, 100) # Lighter gray
                        p_source.alignment = PP_ALIGN.RIGHT # Right align the source

                # --- Add Sources Footer ---
                sources = slide_data.get("sources", [])
                if sources:
                    all_sources.extend(sources)
                    # Add footer text box
                    left = Inches(0.5)
                    top = Inches(7.0) # Bottom of slide
                    width = Inches(9.0)
                    height = Inches(0.5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = "Sources: " + "; ".join(sources)
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(100, 100, 100)

            # --- Add References Slide ---
            if all_sources:
                ref_slide = prs.slides.add_slide(content_layout)
                if ref_slide.shapes.title:
                    ref_slide.shapes.title.text = "References"
                
                # Find body placeholder
                body_shape = None
                for shape in ref_slide.placeholders:
                    if shape.placeholder_format.idx == 10:
                        body_shape = shape
                        break
                if not body_shape: # Fallback
                     for shape in ref_slide.placeholders:
                        if shape.has_text_frame:
                            body_shape = shape
                            break
                
                if body_shape:
                    tf = body_shape.text_frame
                    tf.text = ""
                    unique_sources = sorted(list(set(all_sources)))
                    for source in unique_sources:
                        p = tf.add_paragraph()
                        p.text = source
                        p.level = 0
                        p.font.size = Pt(12)

            prs.save(output_path)
            return ToolResult(output=f"Presentation saved successfully to: {output_path}")

        except Exception as e:
            return ToolResult(error=f"Failed to create PPT: {e}")
